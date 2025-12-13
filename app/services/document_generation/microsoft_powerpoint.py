"""
Serviço para manipulação de Microsoft PowerPoint via Graph API.
Similar ao MicrosoftWordService, mas para apresentações.
"""
from typing import Dict, Any, Optional
import requests
import logging
from .tag_processor import TagProcessor

logger = logging.getLogger(__name__)


class MicrosoftPowerPointService:
    """
    Serviço para manipulação de Microsoft PowerPoint via Graph API.
    """
    
    def __init__(self, credentials: Dict[str, Any]):
        """
        Args:
            credentials: Dict com credenciais Microsoft (access_token, refresh_token, etc.)
                        ou string com access_token (compatibilidade)
        """
        if isinstance(credentials, str):
            # Compatibilidade: aceitar string como access_token
            self.access_token = credentials
        else:
            self.access_token = credentials.get('access_token')
        
        if not self.access_token:
            raise ValueError('access_token não fornecido')
        
        self.base_url = 'https://graph.microsoft.com/v1.0'
        self.headers = {
            'Authorization': f'Bearer {self.access_token}',
            'Content-Type': 'application/json'
        }
    
    def copy_template(self, template_id: str, new_name: str, folder_id: str = None) -> Dict:
        """
        Copia um template do PowerPoint no OneDrive/SharePoint.
        
        Args:
            template_id: ID do arquivo template (driveItem id)
            new_name: Nome da nova apresentação
            folder_id: ID da pasta de destino (opcional)
        
        Returns:
            Dict com id e url da nova apresentação
        """
        # Obter informações do arquivo original
        file_info = requests.get(
            f'{self.base_url}/me/drive/items/{template_id}',
            headers=self.headers
        ).json()
        
        # Determinar pasta de destino
        parent_reference = file_info.get('parentReference', {})
        if folder_id:
            parent_id = folder_id
        else:
            parent_id = parent_reference.get('id')
        
        # Copiar arquivo
        copy_body = {
            'name': new_name
        }
        
        copy_response = requests.post(
            f'{self.base_url}/me/drive/items/{template_id}/copy',
            headers=self.headers,
            json=copy_body
        )
        
        if copy_response.status_code == 202:
            # Copy é assíncrono, aguardar e buscar
            import time
            time.sleep(1)
            
            # Buscar arquivo na pasta
            folder_items = requests.get(
                f'{self.base_url}/me/drive/items/{parent_id}/children',
                headers=self.headers,
                params={'$filter': f"name eq '{new_name}'"}
            ).json()
            
            if folder_items.get('value'):
                new_file = folder_items['value'][0]
                return {
                    'id': new_file['id'],
                    'url': new_file.get('webUrl', f"https://graph.microsoft.com/v1.0/me/drive/items/{new_file['id']}")
                }
        
        copy_response.raise_for_status()
        raise Exception('Falha ao copiar arquivo PowerPoint')
    
    def replace_tags_in_presentation(
        self,
        presentation_id: str,
        data: Dict[str, Any],
        mappings: Dict[str, str] = None
    ) -> None:
        """
        Substitui tags na apresentação PowerPoint.
        
        Nota: Microsoft Graph API não tem endpoint direto para editar conteúdo PowerPoint.
        Usaremos python-pptx para processar o arquivo.
        
        Args:
            presentation_id: ID da apresentação
            data: Dados para substituição
            mappings: Mapeamento de tags para campos
        """
        try:
            # Baixar arquivo
            pptx_content = self._get_presentation_content(presentation_id)
            
            # Processar com python-pptx
            from io import BytesIO
            from pptx import Presentation
            
            prs = Presentation(BytesIO(pptx_content))
            
            # Substituir tags em slides
            for slide in prs.slides:
                # Substituir em shapes de texto
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text = shape.text
                        tags = TagProcessor.extract_tags(text)
                        
                        for tag in tags:
                            field = mappings.get(tag, tag) if mappings else tag
                            value = TagProcessor._get_nested_value(data, field)
                            
                            if value is not None:
                                text = text.replace(f'{{{{{tag}}}}}', str(value))
                            else:
                                text = text.replace(f'{{{{{tag}}}}}', '')
                        
                        shape.text = text
                    
                    # Substituir em tabelas
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text = cell.text
                                tags = TagProcessor.extract_tags(text)
                                
                                for tag in tags:
                                    field = mappings.get(tag, tag) if mappings else tag
                                    value = TagProcessor._get_nested_value(data, field)
                                    
                                    if value is not None:
                                        text = text.replace(f'{{{{{tag}}}}}', str(value))
                                    else:
                                        text = text.replace(f'{{{{{tag}}}}}', '')
                                
                                cell.text = text
            
            # Salvar em buffer
            output = BytesIO()
            prs.save(output)
            output.seek(0)
            
            # Fazer upload do arquivo atualizado
            self._upload_presentation_content(presentation_id, output.read())
            
        except ImportError:
            logger.error('python-pptx não instalado. Instale com: pip install python-pptx')
            raise Exception('Biblioteca python-pptx não disponível')
        except Exception as e:
            logger.exception(f'Erro ao substituir tags no PowerPoint: {str(e)}')
            raise
    
    def _get_presentation_content(self, presentation_id: str) -> bytes:
        """Obtém o conteúdo da apresentação como bytes"""
        response = requests.get(
            f'{self.base_url}/me/drive/items/{presentation_id}/content',
            headers={'Authorization': self.headers['Authorization']}
        )
        response.raise_for_status()
        return response.content
    
    def _upload_presentation_content(self, presentation_id: str, content: bytes) -> None:
        """Faz upload do conteúdo atualizado da apresentação"""
        if len(content) < 4 * 1024 * 1024:  # 4MB
            response = requests.put(
                f'{self.base_url}/me/drive/items/{presentation_id}/content',
                headers={'Authorization': self.headers['Authorization']},
                data=content
            )
            response.raise_for_status()
        else:
            self._upload_large_file(presentation_id, content)
    
    def _upload_large_file(self, presentation_id: str, content: bytes) -> None:
        """Faz upload de arquivo grande usando upload session"""
        session_response = requests.post(
            f'{self.base_url}/me/drive/items/{presentation_id}/createUploadSession',
            headers=self.headers,
            json={
                'item': {
                    '@microsoft.graph.conflictBehavior': 'replace'
                }
            }
        )
        session_response.raise_for_status()
        upload_url = session_response.json()['uploadUrl']
        
        upload_response = requests.put(
            upload_url,
            headers={'Content-Length': str(len(content))},
            data=content
        )
        upload_response.raise_for_status()
    
    def export_as_pdf(self, presentation_id: str) -> bytes:
        """Exporta apresentação PowerPoint como PDF"""
        response = requests.get(
            f'{self.base_url}/me/drive/items/{presentation_id}/content',
            headers={
                'Authorization': self.headers['Authorization'],
                'Accept': 'application/pdf'
            }
        )
        response.raise_for_status()
        return response.content

